"""
build_office_templates.py — 從 sample-data 生成 Office 樣板

產出位置:resources/office-templates/
    01-bar-daily-cases.xlsx            Pattern A 直條 + 7 日 MA
    02-line-three-waves.xlsx           Pattern B 折線
    03-stacked-variants.xlsx           Pattern B 100% 堆疊
    04-stacked-monochrome.xlsx         Pattern E 單色堆疊（重症在底）
    05-pie-age-distribution.xlsx       Pattern B 圓餅（含條件使用註記）
    epidemic-report-template.pptx      嵌入既有範例 PNG 的簡報樣板

執行方式（從 repo 根目錄）：
    python dev-tools/build_office_templates.py

Dev-only 相依（不影響 skill/ runtime）：
    pip install openpyxl python-pptx

成品為靜態樣板：使用者下載即可使用，不需安裝 Python。
"""
from __future__ import annotations

import csv
import sys
from collections import defaultdict
from pathlib import Path

from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.chart.label import DataLabelList
from openpyxl.chart.layout import Layout, ManualLayout
from openpyxl.chart.shapes import GraphicalProperties
from openpyxl.chart.text import RichText
from openpyxl.drawing.colors import ColorChoice
from openpyxl.drawing.fill import ColorChoice as FillColorChoice
from openpyxl.drawing.line import LineProperties
from openpyxl.drawing.text import CharacterProperties, Paragraph, ParagraphProperties, RichTextProperties
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# 把 skill/scripts/ 加進 sys.path 以便 import 規範權威的色票常數
HERE = Path(__file__).resolve().parent
REPO_ROOT = HERE.parent
sys.path.insert(0, str(REPO_ROOT / "skill" / "scripts"))

from epidemic_palette import (  # noqa: E402
    CATEGORICAL,
    LINE_COLORS,
    MONOCHROME,
    PRIMARY,
    PRIMARY_DARK,
    PRIMARY_DARKER,
    PRIMARY_LIGHT,
)

SAMPLE_DATA = REPO_ROOT / "skill" / "assets" / "sample-data"
EXAMPLES = REPO_ROOT / "skill" / "assets" / "examples"
OUTPUT = REPO_ROOT / "resources" / "office-templates"

PAGES_URL = "https://drhao.github.io/epi-dataviz-styleguide/"


# ============== 共用工具 ==============

def hex_no_hash(c: str) -> str:
    return c.lstrip("#").upper()


def read_csv(filename: str) -> list[list[str]]:
    """讀 CSV(處理 UTF-8 BOM),回傳 list of rows"""
    path = SAMPLE_DATA / filename
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return [row for row in csv.reader(f)]


def write_data_sheet(ws, rows: list[list]) -> None:
    """把資料寫進 sheet,第一列加粗作為表頭"""
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row, start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=val)
            if r_idx == 1:
                cell.font = Font(bold=True, color=hex_no_hash(PRIMARY_DARKER))
                cell.fill = PatternFill("solid", fgColor=hex_no_hash("#F2F3F1"))
            cell.alignment = Alignment(horizontal="center" if c_idx > 1 else "left")
    for col_idx in range(1, len(rows[0]) + 1):
        ws.column_dimensions[chr(64 + col_idx)].width = 14


def add_chart_meta(ws, title: str, subtitle: str, pattern_note: str) -> None:
    """在「圖表」sheet 上方寫入標題與引用說明"""
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=16, color=hex_no_hash(PRIMARY_DARKER))
    ws["A2"] = subtitle
    ws["A2"].font = Font(size=11, color=hex_no_hash("#5D675B"))
    ws["A3"] = pattern_note
    ws["A3"].font = Font(size=10, italic=True, color=hex_no_hash(PRIMARY_DARK))
    ws["A4"] = f"完整指引：{PAGES_URL}"
    ws["A4"].font = Font(size=9, color=hex_no_hash("#7A8778"))
    ws.column_dimensions["A"].width = 18


def style_series_fill(series, hex_color: str) -> None:
    """設定 series 填色（直條/堆疊/圓餅用）"""
    series.graphicalProperties = GraphicalProperties(solidFill=hex_no_hash(hex_color))


def style_series_line(series, hex_color: str, width_pt: float = 2.5) -> None:
    """設定 series 線條顏色與寬度（折線用）。width 單位轉 EMU(1pt=12700)"""
    gp = GraphicalProperties()
    gp.line = LineProperties(
        solidFill=hex_no_hash(hex_color),
        w=int(width_pt * 12700),
    )
    series.graphicalProperties = gp


def common_chart_style(chart) -> None:
    """套用本指引通用圖表外觀:無圖例邊框、字體大小、無多餘格線"""
    chart.y_axis.majorGridlines = None  # 移除水平格線(可選,Excel 通常保留)
    chart.x_axis.delete = False
    chart.y_axis.delete = False
    chart.height = 10
    chart.width = 20


# ============== 01 直條圖(Pattern A)==============

def build_bar_daily_cases() -> None:
    """28 天每日新增 + 7 日 trailing MA(本日含前 6 日)"""
    wb = Workbook()
    raw = read_csv("01-daily-cases.csv")

    # 資料 sheet:加入第 4 欄 ma_7(trailing:本日含前 6 日)
    # 前 6 天用自適應窗口(從第 1 天累積到當天),避免斷線
    ws_data = wb.active
    ws_data.title = "資料"
    header = raw[0] + ["ma_7"]
    rows = [header]
    n = len(raw) - 1
    for i, row in enumerate(raw[1:], start=1):
        # i 是 1-based 資料序號,Excel 中對應 row i+1
        excel_row = i + 1
        # Trailing 7 日 MA:窗口為 [max(1, i-6) .. i]
        # 對應 Excel 列:[max(2, excel_row - 6) .. excel_row]
        lo_row = max(2, excel_row - 6)
        ma_formula = f"=AVERAGE(C{lo_row}:C{excel_row})"
        new_row = list(row) + [ma_formula]
        # 轉型數值欄
        new_row[2] = int(row[2])
        rows.append(new_row)
    write_data_sheet(ws_data, rows)

    # 圖表 sheet
    ws_chart = wb.create_sheet("圖表")
    add_chart_meta(
        ws_chart,
        "每日新增確診（28 天）",
        "Pattern A:主色 + 7 日 trailing 移動平均(本日含前 6 日)",
        "色彩:#739A6D 直條 + #374C34 均線。Y 軸從零開始;不截斷座標。前 6 天用自適應累積窗口。",
    )

    # 直條圖
    bar = BarChart()
    bar.type = "col"
    bar.style = 2
    bar.title = None
    bar.y_axis.title = "每日新增（人）"
    bar.x_axis.title = "日期"
    bar.y_axis.scaling.min = 0
    bar.legend.position = "b"

    data_ref = Reference(ws_data, min_col=3, min_row=1, max_row=n + 1, max_col=3)
    cats_ref = Reference(ws_data, min_col=1, min_row=2, max_row=n + 1)
    bar.add_data(data_ref, titles_from_data=True)
    bar.set_categories(cats_ref)
    style_series_fill(bar.series[0], PRIMARY)
    bar.gapWidth = 60  # 對應規範的 barPercentage 0.6 視覺密度
    common_chart_style(bar)

    # 折線:7 日 MA
    line = LineChart()
    ma_ref = Reference(ws_data, min_col=4, min_row=1, max_row=n + 1)
    line.add_data(ma_ref, titles_from_data=True)
    style_series_line(line.series[0], PRIMARY_DARKER, width_pt=2.5)
    line.y_axis.crosses = "autoZero"
    # 把折線疊在直條上(雙軸 chart)
    bar += line

    ws_chart.add_chart(bar, "A6")
    wb.save(OUTPUT / "01-bar-daily-cases.xlsx")


# ============== 02 折線圖(Pattern B)==============

def build_line_three_waves() -> None:
    """三波疫情同期比較,3 條折線"""
    wb = Workbook()
    raw = read_csv("02-weekly-waves.csv")

    ws_data = wb.active
    ws_data.title = "資料"
    rows = [raw[0]]
    for row in raw[1:]:
        rows.append([int(row[0])] + [int(x) for x in row[1:]])
    write_data_sheet(ws_data, rows)
    n = len(raw) - 1

    ws_chart = wb.create_sheet("圖表")
    add_chart_meta(
        ws_chart,
        "三波疫情同期比較",
        "Pattern B:類別配色,最新波次用主色作為焦點",
        "最新一波(2024)用主色 #5D7F58 作為焦點,歷史波次依時間遞遠遞弱(2023 藍 / 2022 黃)。LINE_COLORS 加深版確保 WCAG AA(細線對比 ≥ 3:1)。",
    )

    chart = LineChart()
    chart.title = None
    chart.y_axis.title = "每日新增"
    chart.x_axis.title = "相對天數"
    chart.y_axis.scaling.min = 0
    chart.legend.position = "b"

    data_ref = Reference(ws_data, min_col=2, min_row=1, max_col=4, max_row=n + 1)
    cats_ref = Reference(ws_data, min_col=1, min_row=2, max_row=n + 1)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    # 資料順序為 wave_2022 / 2023 / 2024,顏色反序讓最新一波(2024)用主色
    line_palette = [LINE_COLORS["yellow"], LINE_COLORS["blue"], LINE_COLORS["primary"]]
    for s, color in zip(chart.series, line_palette):
        style_series_line(s, color, width_pt=2.5)

    common_chart_style(chart)
    ws_chart.add_chart(chart, "A6")
    wb.save(OUTPUT / "02-line-three-waves.xlsx")


# ============== 03 堆疊圖(Pattern B 類別)==============

def build_stacked_variants() -> None:
    """變異株消長 100% 堆疊"""
    wb = Workbook()
    raw = read_csv("05-variant-share.csv")

    ws_data = wb.active
    ws_data.title = "資料"
    rows = [raw[0]]
    for row in raw[1:]:
        rows.append([row[0]] + [int(x) for x in row[1:]])
    write_data_sheet(ws_data, rows)
    n = len(raw) - 1
    n_series = len(raw[0]) - 1

    ws_chart = wb.create_sheet("圖表")
    add_chart_meta(
        ws_chart,
        "變異株消長（百分比堆疊)",
        "Pattern B:類別配色 5 色(綠 → 藍 → 黃 → 鴨綠 → 銅)",
        "5 個變異株彼此為獨立類別,使用類別配色。最高優先級(JN.1)用主色。",
    )

    chart = BarChart()
    chart.type = "col"
    chart.grouping = "percentStacked"
    chart.overlap = 100
    chart.title = None
    chart.y_axis.title = "占比"
    chart.x_axis.title = "月份"
    chart.legend.position = "b"

    data_ref = Reference(ws_data, min_col=2, min_row=1,
                         max_col=1 + n_series, max_row=n + 1)
    cats_ref = Reference(ws_data, min_col=1, min_row=2, max_row=n + 1)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    for s, color in zip(chart.series, CATEGORICAL[:n_series]):
        style_series_fill(s, color)

    common_chart_style(chart)
    ws_chart.add_chart(chart, "A6")
    wb.save(OUTPUT / "03-stacked-variants.xlsx")


# ============== 04 堆疊圖(Pattern E 單色,重症在底)==============

def build_stacked_monochrome() -> None:
    """年齡 × 嚴重度 100% 堆疊,單色色階,重症放底部"""
    wb = Workbook()
    raw = read_csv("07-age-severity.csv")
    # 原始:age_group, mild_pct, moderate_pct, severe_pct
    # 為了「深色在底」,序列順序要重排:severe(底) → moderate → mild(頂)
    # MONOCHROME.scale_3 = ["#B4C9B1"(淺), "#739A6D"(中), "#374C34"(深)]
    # 對應:mild=淺、moderate=中、severe=深;堆疊順序 severe → moderate → mild
    reordered_header = [raw[0][0], "severe_pct", "moderate_pct", "mild_pct"]
    reordered_rows = [reordered_header]
    for row in raw[1:]:
        reordered_rows.append([row[0], int(row[3]), int(row[2]), int(row[1])])

    ws_data = wb.active
    ws_data.title = "資料"
    write_data_sheet(ws_data, reordered_rows)
    n = len(reordered_rows) - 1

    ws_chart = wb.create_sheet("圖表")
    add_chart_meta(
        ws_chart,
        "年齡 × 嚴重度（單色堆疊）",
        "Pattern E:單色色階 MONOCHROME.scale_3(淺→中→深)",
        "嚴重度是序數,使用單色色階。重症(最深色 #374C34)放底部作為視覺基底。",
    )

    chart = BarChart()
    chart.type = "col"
    chart.grouping = "percentStacked"
    chart.overlap = 100
    chart.title = None
    chart.y_axis.title = "占比"
    chart.x_axis.title = "年齡組"
    chart.legend.position = "b"

    data_ref = Reference(ws_data, min_col=2, min_row=1,
                         max_col=4, max_row=n + 1)
    cats_ref = Reference(ws_data, min_col=1, min_row=2, max_row=n + 1)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    # 序列順序:severe(底,深) / moderate(中) / mild(頂,淺)
    severity_colors = [
        MONOCHROME["scale_3"][2],  # severe → #374C34
        MONOCHROME["scale_3"][1],  # moderate → #739A6D
        MONOCHROME["scale_3"][0],  # mild → #B4C9B1
    ]
    for s, color in zip(chart.series, severity_colors):
        style_series_fill(s, color)

    common_chart_style(chart)
    ws_chart.add_chart(chart, "A6")
    wb.save(OUTPUT / "04-stacked-monochrome.xlsx")


# ============== 05 圓餅圖(Pattern B + 條件使用註記)==============

def build_pie_age_distribution() -> None:
    """年齡分布 — 聚合為 5 組以符合圓餅圖規範(類別 ≤ 5)"""
    wb = Workbook()
    raw = read_csv("10-age-gender.csv")
    # 原資料 9 個年齡組,聚合為 5 組
    buckets = {
        "0-19": ["0-9", "10-19"],
        "20-39": ["20-29", "30-39"],
        "40-59": ["40-49", "50-59"],
        "60-79": ["60-69", "70-79"],
        "80+": ["80+"],
    }
    age_totals: dict[str, int] = defaultdict(int)
    by_age = {row[0]: int(row[1]) + int(row[2]) for row in raw[1:]}
    for label, members in buckets.items():
        for m in members:
            age_totals[label] += by_age.get(m, 0)

    rows = [["age_bucket", "population_thousand"]]
    for label in buckets:
        rows.append([label, age_totals[label]])

    ws_data = wb.active
    ws_data.title = "資料"
    write_data_sheet(ws_data, rows)

    ws_chart = wb.create_sheet("圖表")
    add_chart_meta(
        ws_chart,
        "人口年齡分布(聚合為 5 組)",
        "Pattern B:類別配色;原資料 9 組已聚合為 5 組,符合圓餅圖規範",
        "圓餅圖條件使用:類別 ≤ 5 且傳達占比才用,否則改用直條圖。本範例已將 9 組聚合為 5 組。",
    )

    chart = PieChart()
    chart.title = None
    chart.legend.position = "r"
    chart.dataLabels = DataLabelList(showPercent=True)

    data_ref = Reference(ws_data, min_col=2, min_row=1, max_row=6)
    cats_ref = Reference(ws_data, min_col=1, min_row=2, max_row=6)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    # 圓餅圖只有一個 series,要為每個 data point 著色
    from openpyxl.chart.marker import DataPoint
    pts = []
    for idx, color in enumerate(CATEGORICAL[:5]):
        dp = DataPoint(idx=idx)
        dp.graphicalProperties = GraphicalProperties(solidFill=hex_no_hash(color))
        pts.append(dp)
    chart.series[0].data_points = pts

    chart.height = 10
    chart.width = 18
    ws_chart.add_chart(chart, "A6")
    wb.save(OUTPUT / "05-pie-age-distribution.xlsx")


# ============== PPT 樣板 ==============

def build_pptx_template() -> None:
    """6 張投影片的疫情週報範例,嵌入既有 PNG"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def add_title_bar(slide, title_text: str, subtitle_text: str = "") -> None:
        """頂部主色標題列"""
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
            prs.slide_width, Inches(0.8),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0x73, 0x9A, 0x6D)
        bar.line.fill.background()
        bar.shadow.inherit = False

        tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.12),
                                       prs.slide_width - Inches(0.8), Inches(0.6))
        tf = tx.text_frame
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        if subtitle_text:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle_text
            r2.font.size = Pt(11)
            r2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xE7)

    def add_caption(slide, text: str, y: float) -> None:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(y),
                                       prs.slide_width - Inches(1.2), Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0x44, 0x4C, 0x43)

    # --- Slide 1:封面 ---
    s1 = prs.slides.add_slide(blank_layout)
    bar = s1.shapes.add_shape(
        __import__("pptx.enum.shapes", fromlist=["MSO_SHAPE"]).MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(2.8), prs.slide_width, Inches(2.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x73, 0x9A, 0x6D)
    bar.line.fill.background()
    bar.shadow.inherit = False

    tx = s1.shapes.add_textbox(Inches(0.8), Inches(3.0),
                                prs.slide_width - Inches(1.6), Inches(1.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "疫情週報範例"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Epidemic Weekly Report Template"
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xE7)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = f"依據:{PAGES_URL}"
    r3.font.size = Pt(12)
    r3.font.color.rgb = RGBColor(0xE8, 0xEE, 0xE7)

    # --- Slide 2-5:嵌入 PNG ---
    embed_slides = [
        {
            "title": "每日新增確診",
            "subtitle": "Pattern A:主色直條 + 7 日 trailing 移動平均",
            "png": "01b-bar-daily-with-ma.png",
            "caption": "每日填報受週末效應影響,以 7 日 trailing 均線(本日含前 6 日)呈現潛在趨勢。Y 軸從零開始。",
        },
        {
            "title": "今年 vs 去年同期(含歷史範圍)",
            "subtitle": "Pattern A:主色焦點 + 中性灰歷史範圍帶",
            "png": "02c-line-year-over-year.png",
            "caption": "今年用主色 #5D7F58 作為焦點,去年同期用中性灰虛線,歷史範圍(±1 SD)用灰色填充帶,讀者能立即看出當前是否偏離常態。",
        },
        {
            "title": "變異株消長",
            "subtitle": "Pattern B:5 色類別配色 100% 堆疊",
            "png": "04a-stacked-100-percent.png",
            "caption": "5 個變異株彼此為獨立類別,使用類別配色。主流變異株(JN.1)用主色。",
        },
        {
            "title": "年齡 × 嚴重度(單色)",
            "subtitle": "Pattern E:單色色階,重症(深色)放底部",
            "png": "10a-mono-stacked-severity.png",
            "caption": "嚴重度為序數資料,使用單色色階傳達層次。最深色置於堆疊底部作為視覺基底。",
        },
    ]
    for spec in embed_slides:
        slide = prs.slides.add_slide(blank_layout)
        add_title_bar(slide, spec["title"], spec["subtitle"])
        png_path = EXAMPLES / spec["png"]
        if not png_path.exists():
            print(f"  ⚠ 缺少範例 PNG:{png_path}")
            continue
        # 圖片置中放在中段
        pic = slide.shapes.add_picture(
            str(png_path),
            Inches(1.5), Inches(1.1),
            height=Inches(5.2),
        )
        # 微調水平置中
        pic.left = int((prs.slide_width - pic.width) / 2)
        add_caption(slide, spec["caption"], y=6.6)

    # --- Slide 6:色票與引用 ---
    s6 = prs.slides.add_slide(blank_layout)
    add_title_bar(s6, "色票與引用", "依本指引使用 — 主色 #739A6D + 類別配色 + 強調紅僅限警示")

    swatches = [
        ("#739A6D", "主色 Sage Green"),
        ("#587A9D", "Slate Blue"),
        ("#C8A041", "Mustard"),
        ("#49888D", "Teal"),
        ("#916E46", "Bronze"),
        ("#955F71", "Plum"),
        ("#BE373C", "Alert Red(僅警示)"),
    ]
    from pptx.enum.shapes import MSO_SHAPE
    swatch_w, swatch_h = Inches(1.4), Inches(1.4)
    start_x, y = Inches(0.8), Inches(1.4)
    gap = Inches(0.2)
    for i, (hex_c, label) in enumerate(swatches):
        x = start_x + (swatch_w + gap) * i
        sh = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, swatch_w, swatch_h)
        sh.fill.solid()
        rgb = tuple(int(hex_c[j:j+2], 16) for j in (1, 3, 5))
        sh.fill.fore_color.rgb = RGBColor(*rgb)
        sh.line.fill.background()
        sh.shadow.inherit = False
        # 標籤
        lbl = s6.shapes.add_textbox(x, y + swatch_h + Inches(0.05),
                                     swatch_w, Inches(0.7))
        tf = lbl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hex_c
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x37, 0x4C, 0x34)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(9)
        r2.font.color.rgb = RGBColor(0x5D, 0x67, 0x5B)

    # 底部引用區
    ref_box = s6.shapes.add_textbox(Inches(0.8), Inches(4.4),
                                     prs.slide_width - Inches(1.6), Inches(2.5))
    rf = ref_box.text_frame
    rf.word_wrap = True
    for line, bold, size, color in [
        ("4 項核心原則", True, 16, (0x37, 0x4C, 0x34)),
        ("清晰優先 / 誠實呈現 / 一致性 / 負責任溝通", False, 12, (0x44, 0x4C, 0x43)),
        ("", False, 8, (0xFF, 0xFF, 0xFF)),
        ("使用本樣板時請遵守:", True, 12, (0x37, 0x4C, 0x34)),
        ("• 直條圖 Y 軸從零開始,不截斷座標(折線視情境)", False, 11, (0x44, 0x4C, 0x43)),
        ("• 紅色(#BE373C)僅用於警示,不作一般類別色", False, 11, (0x44, 0x4C, 0x43)),
        ("• 7 日移動平均使用 trailing(本日含前 6 日)", False, 11, (0x44, 0x4C, 0x43)),
        ("• 單色色階堆疊時,最深色放底部", False, 11, (0x44, 0x4C, 0x43)),
        ("", False, 8, (0xFF, 0xFF, 0xFF)),
        (f"完整指引:{PAGES_URL}", False, 11, (0x73, 0x9A, 0x6D)),
    ]:
        if rf.paragraphs[0].text == "" and len(rf.paragraphs) == 1 and line:
            p = rf.paragraphs[0]
        else:
            p = rf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)

    prs.save(OUTPUT / "epidemic-report-template.pptx")


# ============== 樣板說明 README ==============

README_CONTENT = """# Office 圖表樣板

從本指引 sample-data 自動生成的 Excel / PowerPoint 樣板。**下載即用,不需 Python 環境。**

## 包含內容

| 檔案 | 圖表類型 | 配色模式 | 資料來源 |
|------|---------|---------|---------|
| `01-bar-daily-cases.xlsx` | 直條 + 7 日 MA | Pattern A(主色 + 強調) | 28 天每日新增(虛構) |
| `02-line-three-waves.xlsx` | 折線(3 條) | Pattern B(類別配色) | 三波疫情同期比較 |
| `03-stacked-variants.xlsx` | 100% 堆疊長條 | Pattern B(類別配色) | 變異株消長 6 個月 |
| `04-stacked-monochrome.xlsx` | 100% 堆疊長條 | **Pattern E(單色,重症在底)** | 年齡 × 嚴重度 |
| `05-pie-age-distribution.xlsx` | 圓餅 | Pattern B + 條件使用註記 | 年齡分布(聚合為 5 組) |
| `epidemic-report-template.pptx` | 6 頁簡報樣板 | — | 嵌入既有範例 PNG |

## 使用方式

### Excel 樣板

1. 下載對應的 `.xlsx`
2. 切到「資料」分頁,改成你的真實資料(欄位結構與表頭請保持不變)
3. 切到「圖表」分頁,圖表自動連動更新
4. 「圖表」分頁頂部的引用註記建議保留,提醒讀者本圖遵循疫情視覺化指引

### PowerPoint 樣板

1. 下載 `epidemic-report-template.pptx`
2. 6 張投影片:封面 + 4 張嵌入範例圖 + 色票/原則摘要
3. 把嵌入的 PNG 換成你自己的圖表(右鍵 → 變更圖片)
4. 標題列主色 `#739A6D` 已套用,可直接調整文字

## 重要原則(套用樣板時請遵守)

- **不要修改主色** `#739A6D` — 這是組織色彩識別
- **紅色 `#BE373C` 僅用於警示**,不可作一般類別色
- **直條圖 Y 軸從零開始**,不截斷座標軸誤導比例(折線視情境)
- **7 日移動平均**使用 trailing(本日含前 6 日,即 i-6 到 i),前 6 天用自適應累積窗口
- **單色色階堆疊**最深色放底部(本範例 04 已示範:重症在底)

## 重新生成

樣板由 `dev-tools/build_office_templates.py` 從 `skill/assets/sample-data/` 生成。色票若調整,重跑:

```bash
pip install openpyxl python-pptx
python3 dev-tools/build_office_templates.py
```

詳見 [`dev-tools/README.md`](../../dev-tools/README.md)。

## 完整指引

線上版:https://drhao.github.io/epi-dataviz-styleguide/
"""


def write_readme() -> None:
    (OUTPUT / "README.md").write_text(README_CONTENT, encoding="utf-8")


# ============== 主流程 ==============

def main() -> int:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    print("=" * 60)
    print("生成 Office 樣板")
    print("=" * 60)

    builders = [
        ("01-bar-daily-cases.xlsx", build_bar_daily_cases),
        ("02-line-three-waves.xlsx", build_line_three_waves),
        ("03-stacked-variants.xlsx", build_stacked_variants),
        ("04-stacked-monochrome.xlsx", build_stacked_monochrome),
        ("05-pie-age-distribution.xlsx", build_pie_age_distribution),
        ("epidemic-report-template.pptx", build_pptx_template),
    ]
    for name, fn in builders:
        print(f"  → {name}")
        fn()

    write_readme()
    print(f"  → README.md")
    print()
    print(f"✓ 完成,輸出位置:{OUTPUT.relative_to(REPO_ROOT)}/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
